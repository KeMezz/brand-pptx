"""brand-pptx — theme.json を読み込み、ブランド準拠スライドを組み立てる薄いヘルパー。

SKILL.md にインライン展開しているロジックの実体。make_sample.py から利用する。
"""
import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

_RELS_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def find_skill_dir() -> Path:
    """theme.json と同じ階層（スキルのルート）を探す。"""
    cands = []
    pr = os.environ.get("CLAUDE_PLUGIN_ROOT")
    if pr:
        cands.append(Path(pr))
    cands += [
        Path.home() / ".claude" / "skills" / "brand-pptx",
        Path.home() / ".agents" / "skills" / "brand-pptx",
        Path.home() / ".config" / "skills" / "brand-pptx",
        Path(__file__).resolve().parent.parent,
        Path.cwd(),
    ]
    for c in cands:
        if (c / "theme.json").exists():
            return c
    raise FileNotFoundError("theme.json が見つかりません。")


class Brand:
    """theme.json をロードし、色/フォント/サイズと生成ヘルパーを提供する。"""

    def __init__(self, skill_dir: Path | None = None):
        self.dir = skill_dir or find_skill_dir()
        self.theme = json.loads((self.dir / "theme.json").read_text(encoding="utf-8"))
        self.C = {k: RGBColor.from_string(v) for k, v in self.theme["colors"].items()}
        self.F = self.theme["fonts"]
        self.TS = self.theme["typeScale"]
        self.LM = self.theme["layoutMap"]
        # opt-in: layoutMap の役割が null のとき、テンプレートを変更せず theme 色で合成する
        self.draw_missing = bool((self.theme.get("setup") or {}).get("drawMissingRoles", False))

    # --- プレゼン/テンプレート ---
    def new_presentation(self):
        tpath = (self.theme.get("template") or {}).get("path")
        if tpath:
            p = Path(tpath)
            if not p.is_absolute():
                p = self.dir / tpath
            prs = Presentation(str(p))
            self.using_custom = True
        else:
            prs = Presentation()
            self.using_custom = False
        return prs

    @staticmethod
    def delete_all_slides(prs):
        while len(prs.slides._sldIdLst):
            rId = prs.slides._sldIdLst[0].attrib[_RELS_ID]
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    @staticmethod
    def ph(slide, idx):
        for p in slide.placeholders:
            if p.placeholder_format.idx == idx:
                return p
        return None

    @staticmethod
    def title_ph(slide):
        """タイトルプレースホルダーを返す。idx=0 を優先し、無ければ型で探す。

        自社テンプレートでは title が idx=0 とは限らないため、
        TITLE / CENTER_TITLE 型へフォールバックする。
        """
        for p in slide.placeholders:
            if p.placeholder_format.idx == 0:
                return p
        for p in slide.placeholders:
            if p.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return p
        return None

    def round_corners(self, shape):
        """ROUNDED_RECTANGLE の角丸半径を theme.style.cornerRadiusIn に合わせる。"""
        r_in = (self.theme.get("style") or {}).get("cornerRadiusIn")
        if r_in is None:
            return
        short = min(int(shape.width), int(shape.height))
        if short <= 0:
            return
        adj = max(0.0, min(0.5, int(Inches(r_in)) / short))
        try:
            shape.adjustments[0] = adj
        except Exception:
            pass

    @staticmethod
    def _flat(shape):
        """テーマ由来のスタイル（プリセットの影など）を外してフラットにする。"""
        st = shape._element.find(qn("p:style"))
        if st is not None:
            shape._element.remove(st)

    @staticmethod
    def set_margins(tf, l=Inches(0.1), r=Inches(0.1), t=Inches(0.05), b=Inches(0.05)):
        tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = l, r, t, b
        tf.word_wrap = True

    def style_text(self, tf, size, color, bold=False, align=PP_ALIGN.LEFT,
                   font=None, anchor=None):
        if anchor is not None:
            tf.vertical_anchor = anchor
        for para in tf.paragraphs:
            para.alignment = align
            runs = para.runs or [para.add_run()]
            for r in runs:
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.color.rgb = color
                if font:
                    r.font.name = font

    # --- 不足ロールの合成（opt-in: theme.setup.drawMissingRoles）---
    def _dims(self, prs):
        m = Inches((self.theme.get("style") or {}).get("pageMarginIn", 0.6))
        return prs.slide_width, prs.slide_height, m

    @staticmethod
    def _blank_layout(prs):
        # プレースホルダーが最も少ない＝最も白紙に近いレイアウトを土台にする
        return min(prs.slide_layouts, key=lambda L: len(list(L.placeholders)))

    def _role_layout(self, prs, role):
        """layoutMap[role] が int ならそのレイアウト、null/未設定なら白紙を返す。
        戻り値: (slide, synthesized)"""
        idx = self.LM.get(role)
        if idx is not None:
            return prs.slides.add_slide(prs.slide_layouts[idx]), False
        return prs.slides.add_slide(self._blank_layout(prs)), True

    def _rect(self, slide, shp, x, y, w, h, color):
        s = slide.shapes.add_shape(shp, x, y, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
        self._flat(s)
        return s

    def _synth_chrome(self, prs, slide, role):
        """flag が on のとき、テンプレートに無いロールの装飾を theme 色で「スライドに」描く
        （テンプレート/マスターは一切変更しない）。"""
        if not self.draw_missing:
            return
        W, H, m = self._dims(prs)
        if role in ("section", "ending"):
            self._rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, H, self.C["primary"])
            if role == "section":
                bh = Inches(1.2)
                self._rect(slide, MSO_SHAPE.RECTANGLE, m, (H - bh) // 2, Inches(0.14), bh, self.C["accent"])
        else:  # cover / content / contentVisual
            sq = Inches(0.32)
            sy = Inches(2.0) if role == "cover" else Inches(0.62)
            self._rect(slide, MSO_SHAPE.RECTANGLE, m, sy, sq, sq, self.C["accent"])
            if role in ("content", "contentVisual"):
                self._rect(slide, MSO_SHAPE.RECTANGLE, m, H - Inches(0.55), W - m * 2, Inches(0.025), self.C["border"])

    def _synth_text(self, prs, slide, role, title, subtitle=None):
        """合成スライドにタイトル（必要ならサブタイトル）をテキストボックスで描く。"""
        W, H, m = self._dims(prs)
        dark_bg = self.draw_missing and role in ("section", "ending")
        color = self.C["white"] if dark_bg else self.C["dark"]
        if role == "cover":
            size, y, h, align, tx = self.TS["coverTitle"], Inches(2.35), Inches(1.6), PP_ALIGN.LEFT, m
        elif role == "section":
            size, h, align, tx = self.TS["sectionTitle"], Inches(1.7), PP_ALIGN.LEFT, m + Inches(0.45)
            y = (H - h) // 2
        elif role == "ending":
            size, h, align, tx = self.TS["sectionTitle"], Inches(1.4), PP_ALIGN.CENTER, m
            y = (H - h) // 2
        else:  # content / contentVisual
            size, y, h, align, tx = self.TS["title"], Inches(0.55), Inches(0.9), PP_ALIGN.LEFT, m
        box = slide.shapes.add_textbox(tx, y, W - tx - m, h)
        box.text_frame.paragraphs[0].text = title
        self.style_text(box.text_frame, size, color, bold=True, align=align,
                        font=self.F["headingCJK"], anchor=MSO_ANCHOR.MIDDLE)
        if subtitle and role == "cover":
            sb = slide.shapes.add_textbox(tx, y + h + Inches(0.05), W - tx - m, Inches(0.6))
            sb.text_frame.paragraphs[0].text = subtitle
            self.style_text(sb.text_frame, self.TS["coverSubtitle"], self.C["gray"], font=self.F["bodyCJK"])

    # --- スライドビルダー ---
    def add_cover(self, prs, title, subtitle=None):
        slide, synth = self._role_layout(prs, "cover")
        if synth:
            self._synth_chrome(prs, slide, "cover")
            self._synth_text(prs, slide, "cover", title, subtitle)
            return slide
        t = self.title_ph(slide)
        if t:
            t.text = title
            self.style_text(t.text_frame, self.TS["coverTitle"], self.C["dark"],
                            bold=True, font=self.F["headingCJK"], anchor=MSO_ANCHOR.MIDDLE)
        if subtitle:
            s = self.ph(slide, 1)
            if s:
                s.text = subtitle
                self.style_text(s.text_frame, self.TS["coverSubtitle"], self.C["gray"],
                                font=self.F["bodyCJK"])
        return slide

    def add_section(self, prs, title):
        # 背景・装飾はテンプレートのレイアウトが持つ。ここでは文字を流し込むだけ。
        slide, synth = self._role_layout(prs, "section")
        if synth:
            self._synth_chrome(prs, slide, "section")
            self._synth_text(prs, slide, "section", title)
            return slide
        t = self.title_ph(slide)
        if t:
            t.text = title
            self.style_text(t.text_frame, self.TS["sectionTitle"], self.C["white"],
                            bold=True, font=self.F["headingCJK"], anchor=MSO_ANCHOR.MIDDLE)
        return slide

    def add_content(self, prs, title, visual=True):
        key = "contentVisual" if visual else "content"
        slide, synth = self._role_layout(prs, key)
        if synth:
            self._synth_chrome(prs, slide, key)
            self._synth_text(prs, slide, key, title)
            return slide
        t = self.title_ph(slide)
        if t:
            t.text = title
            self.style_text(t.text_frame, self.TS["title"], self.C["dark"],
                            bold=True, font=self.F["headingCJK"])
        return slide

    def add_ending(self, prs, title="ご清聴ありがとうございました"):
        slide, synth = self._role_layout(prs, "ending")
        if synth:
            self._synth_chrome(prs, slide, "ending")
            self._synth_text(prs, slide, "ending", title)
            return slide
        t = self.title_ph(slide)
        if t:
            t.text = title
            self.style_text(t.text_frame, self.TS["sectionTitle"], self.C["white"],
                            bold=True, align=PP_ALIGN.CENTER, font=self.F["headingCJK"],
                            anchor=MSO_ANCHOR.MIDDLE)
        return slide

    # --- 視覚パーツ ---
    def add_kpi_card(self, slide, x, y, w, h, value, label):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid(); card.fill.fore_color.rgb = self.C["bg"]; card.line.fill.background()
        self.round_corners(card); self._flat(card)
        val_h = int(h * 0.55)
        b1 = slide.shapes.add_textbox(x, y, w, val_h); tf = b1.text_frame
        self.set_margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].text = str(value)
        tf.paragraphs[0].font.size = Pt(self.TS["kpiValue"]); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.C["accent"]; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        b2 = slide.shapes.add_textbox(x, y + val_h, w, h - val_h); tf2 = b2.text_frame
        self.set_margins(tf2); tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf2.paragraphs[0].text = label
        tf2.paragraphs[0].font.size = Pt(self.TS["caption"]); tf2.paragraphs[0].font.color.rgb = self.C["gray"]
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_kpi_row(self, slide, items, x, y, total_w, card_h, gap=None):
        """items: [(value, label), ...] を total_w に等幅で敷き詰める（左右いっぱい＝バランス）。"""
        n = len(items)
        if gap is None:
            gap = Inches((self.theme.get("style") or {}).get("blockGapIn", 0.4))
        card_w = int((total_w - gap * (n - 1)) / n)
        for i, (value, label) in enumerate(items):
            self.add_kpi_card(slide, x + (card_w + gap) * i, y, card_w, card_h, value, label)

    def add_progress_bar(self, slide, x, y, w, h, progress):
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        bg.fill.solid(); bg.fill.fore_color.rgb = self.C["bg"]; bg.line.fill.background()
        self.round_corners(bg); self._flat(bg)
        bw = int(w * max(0.0, min(1.0, progress)))
        if bw > 0:
            bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, h)
            bar.fill.solid(); bar.fill.fore_color.rgb = self.C["accent"]; bar.line.fill.background()
            self.round_corners(bar); self._flat(bar)

    def add_step_flow(self, slide, steps, x, y, total_w):
        n = len(steps); circle = Inches(0.6); step_w = int(total_w / n)
        for i, text in enumerate(steps):
            cx = x + step_w * i + (step_w - circle) // 2
            c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, y, circle, circle)
            c.fill.solid(); c.fill.fore_color.rgb = self.C["accent"]; c.line.fill.background()
            self._flat(c)
            nb = slide.shapes.add_textbox(cx, y, circle, circle); tf = nb.text_frame
            self.set_margins(tf, Inches(0), Inches(0), Inches(0), Inches(0)); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.paragraphs[0].text = str(i + 1)
            tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self.C["white"]; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            lb = slide.shapes.add_textbox(x + step_w * i, y + circle + Inches(0.1), step_w, Inches(0.5))
            tf2 = lb.text_frame; self.set_margins(tf2)
            tf2.paragraphs[0].text = text
            tf2.paragraphs[0].font.size = Pt(self.TS["caption"]); tf2.paragraphs[0].font.color.rgb = self.C["dark"]
            tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
            if i < n - 1:
                # 矢印は「今の円の右端」と「次の円の左端」の中間に置く
                cx_next = x + step_w * (i + 1) + (step_w - circle) // 2
                arrow_w = Inches(0.4)
                ax = (cx + circle) + ((cx_next - (cx + circle)) - arrow_w) // 2
                ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax,
                                            y + Inches(0.15), arrow_w, Inches(0.3))
                ar.fill.solid(); ar.fill.fore_color.rgb = self.C["accentLight"]; ar.line.fill.background()
                self._flat(ar)

    def add_data_table(self, slide, x, y, headers, rows, w):
        n_rows, n_cols = len(rows) + 1, len(headers)
        table = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(0.4 * n_rows)).table
        table.first_row = False      # 既定スタイルのヘッダー装飾を無効化（自前で塗る）
        table.horz_banding = False   # 行の縞模様を無効化
        for j, h in enumerate(headers):
            cell = table.cell(0, j); cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = self.C["bg"]
            self.style_text(cell.text_frame, self.TS["caption"], self.C["dark"],
                            bold=True, font=self.F["bodyCJK"])
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j); cell.text = str(val)
                self.style_text(cell.text_frame, self.TS["caption"], self.C["dark"],
                                font=self.F["bodyCJK"])
