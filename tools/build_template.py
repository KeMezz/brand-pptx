"""theme.json の配色を読み込み、オリジナルデザインの既定テンプレート
（assets/template.pptx, 16:9）を生成する。

各レイアウトに装飾（カバーの右パネル、セクションの全面背景、コンテンツの
アクセント＋フッター罫線、エンディングの全面背景）を「レイアウト自体」へ
焼き込むので、生成時はプレースホルダーを埋めるだけで装飾が付く。

    python3 tools/build_template.py

theme.json の colors を変えてから再実行すれば、既定テンプレートの配色も変わる。
"""
import json
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
theme = json.loads((ROOT / "theme.json").read_text(encoding="utf-8"))
C = {k: RGBColor.from_string(v) for k, v in theme["colors"].items()}

EMU = 914400
W, H = Inches(13.333), Inches(7.5)  # 16:9


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 装飾を組み立てる作業用スライド（最後に全削除）
    scratch = prs.slides.add_slide(prs.slide_layouts[6])

    def rect(shape_type, x, y, w, h, color):
        sp = scratch.shapes.add_shape(shape_type, x, y, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        sp.shadow.inherit = False
        # テーマ由来のスタイル（プリセットの影など）を外してフラットにする
        st = sp._element.find(qn("p:style"))
        if st is not None:
            sp._element.remove(st)
        return sp

    def inject(layout, shapes, start=2):
        """作業用スライドで作った図形を layout の spTree に背面側から挿入する。"""
        spTree = layout.shapes._spTree
        pos = start
        for sp in shapes:
            spTree.insert(pos, deepcopy(sp._element))
            pos += 1
        for sp in shapes:  # 作業用スライドからは取り除く
            sp._element.getparent().remove(sp._element)

    def place(layout, idx, x, y, w, h):
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == idx:
                ph.left, ph.top, ph.width, ph.height = x, y, w, h
                return ph
        return None

    def drop_ph(layout, idxs):
        for ph in list(layout.placeholders):
            if ph.placeholder_format.idx in idxs:
                ph._element.getparent().remove(ph._element)

    # ---- cover (layout 0: Title Slide) ----
    cov = prs.slide_layouts[0]
    inject(cov, [
        rect(MSO_SHAPE.RECTANGLE, Inches(9.0), 0, Inches(4.333), H, C["primary"]),      # 右パネル
        rect(MSO_SHAPE.OVAL, Inches(8.25), Inches(2.35), Inches(1.5), Inches(1.5), C["accent"]),  # アクセント円
        rect(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.95), Inches(0.32), Inches(0.32), C["accent"]),  # 角アクセント
    ])
    place(cov, 0, Inches(0.9), Inches(2.4), Inches(7.2), Inches(1.6))
    place(cov, 1, Inches(0.9), Inches(4.15), Inches(7.0), Inches(0.8))

    # ---- section (layout 2: Section Header) ----
    sec = prs.slide_layouts[2]
    inject(sec, [
        rect(MSO_SHAPE.RECTANGLE, 0, 0, W, H, C["primary"]),                            # 全面背景
        rect(MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(3.15), Inches(0.14), Inches(1.2), C["accent"]),  # アクセントバー（垂直中央 3.75"）
    ])
    place(sec, 0, Inches(1.5), Inches(2.9), Inches(10.3), Inches(1.7))   # title（中央 3.75"・テキストは MIDDLE アンカー）
    place(sec, 1, Inches(1.5), Inches(2.25), Inches(10.3), Inches(0.6))  # eyebrow（任意・未使用）

    # ---- content (layout 1: Title and Content) ----
    con = prs.slide_layouts[1]
    inject(con, [
        rect(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.62), Inches(0.32), Inches(0.32), C["accent"]),  # アクセント角
        rect(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.95), Inches(11.73), Inches(0.025), C["border"]),  # フッター罫線
    ])
    place(con, 0, Inches(1.3), Inches(0.55), Inches(11.0), Inches(0.75))
    place(con, 1, Inches(0.85), Inches(1.65), Inches(11.6), Inches(5.0))

    # ---- contentVisual (layout 5: Title Only) ----
    cv = prs.slide_layouts[5]
    inject(cv, [
        rect(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.62), Inches(0.32), Inches(0.32), C["accent"]),
        rect(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.95), Inches(11.73), Inches(0.025), C["border"]),
    ])
    place(cv, 0, Inches(1.3), Inches(0.55), Inches(11.0), Inches(0.75))

    # ---- ending (layout 7: Content with Caption) ----
    end = prs.slide_layouts[7]
    drop_ph(end, {1, 2})  # 余分なプレースホルダーを除去（title だけ使う）
    inject(end, [
        rect(MSO_SHAPE.RECTANGLE, 0, 0, W, H, C["primary"]),
        rect(MSO_SHAPE.RECTANGLE, Inches(6.49), Inches(2.55), Inches(0.34), Inches(0.34), C["accent"]),
    ])
    place(end, 0, Inches(1.5), Inches(3.05), Inches(10.33), Inches(1.4))

    # 作業用スライドを削除（テンプレートはスライド 0 枚で出荷）
    rid_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    while len(prs.slides._sldIdLst):
        rId = prs.slides._sldIdLst[0].get(rid_ns)
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    out = ROOT / "assets" / "template.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"saved: {out}  (16:9, layouts decorated: cover=0 section=2 content=1 contentVisual=5 ending=7)")


if __name__ == "__main__":
    sys.exit(main())
