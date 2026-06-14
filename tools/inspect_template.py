"""テンプレート PPTX のスライドレイアウト一覧とプレースホルダーを表示する。

    python3 tools/inspect_template.py                 # 内蔵既定テンプレート
    python3 tools/inspect_template.py path/to/your.pptx

出力された index を theme.json の layoutMap（cover / section / content /
contentVisual / ending → どの index か）に書き写すと、自社テンプレートに対応できる。
"""
import sys
from pptx import Presentation


def main(path=None):
    prs = Presentation(path) if path else Presentation()
    print(f"template: {path or 'builtin-default'}")
    print(f"slide size: {prs.slide_width / 914400:.3f}in x {prs.slide_height / 914400:.3f}in\n")
    for i, layout in enumerate(prs.slide_layouts):
        phs = ", ".join(
            f"idx={p.placeholder_format.idx}:{p.placeholder_format.type}"
            for p in layout.placeholders
        )
        print(f"[{i:>2}] {layout.name}")
        print(f"      placeholders: {phs or '(none)'}")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else None)
